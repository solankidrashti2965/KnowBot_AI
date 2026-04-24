"""
Document Parser — extracts and chunks text from various file formats.
Supports: PDF, DOCX, PPTX, TXT, MD, and common Images.
"""

import os
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
import io
from langchain_text_splitters import RecursiveCharacterTextSplitter
from typing import List, Tuple, Dict


def extract_text_from_file(file_path: str) -> Tuple[List[Dict], int]:
    """
    Extract and chunk text from a file based on its extension.

    Returns:
    chunks     — list of {"content": str, "page": int}
    page_count — total 'pages' or 'slides' in the document
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == ".pdf":
        return _parse_pdf(file_path)
    elif ext in [".docx", ".doc"]:
        return _parse_docx(file_path)
    elif ext in [".pptx", ".ppt"]:
        return _parse_pptx(file_path)
    elif ext in [".txt", ".md", ".py", ".js", ".json", ".csv", ".html"]:
        return _parse_text(file_path)
    elif ext in [".png", ".jpg", ".jpeg", ".webp"]:
        return _parse_image(file_path)
    else:
        raise ValueError(f"Unsupported file format: {ext}")


def _parse_pdf(file_path: str) -> Tuple[List[Dict], int]:
    doc = fitz.open(file_path)
    page_count = len(doc)
    pages_text = []

    for page_num, page in enumerate(doc):
        text = page.get_text("text")
        if text.strip():
            pages_text.append({"page": page_num + 1, "text": text.strip()})
    
    doc.close()
    return _chunk_text(pages_text, page_count)


def _parse_docx(file_path: str) -> Tuple[List[Dict], int]:
    doc = DocxDocument(file_path)
    full_text = []
    for para in doc.paragraphs:
        if para.text.strip():
            full_text.append(para.text.strip())
    
    # Word doesn't have reliable "page numbers" in raw text, treat as 1 group
    pages_text = [{"page": 1, "text": "\n".join(full_text)}]
    return _chunk_text(pages_text, 1)


def _parse_pptx(file_path: str) -> Tuple[List[Dict], int]:
    prs = Presentation(file_path)
    slides_text = []
    
    for i, slide in enumerate(prs.slides):
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())
        
        if text_parts:
            slides_text.append({"page": i + 1, "text": "\n".join(text_parts)})
            
    return _chunk_text(slides_text, len(prs.slides))


def _parse_text(file_path: str) -> Tuple[List[Dict], int]:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    
    pages_text = [{"page": 1, "text": text.strip()}]
    return _chunk_text(pages_text, 1)


def _parse_image(file_path: str) -> Tuple[List[Dict], int]:
    """
    Standalone image parsing.
    For high-quality OCR, Tesseract should be used.
    For now, we use PyMuPDF's built-in OCR if available or just return image metadata.
    """
    try:
        # Try PyMuPDF's OCR or simple text extraction if it's an image-pdf
        doc = fitz.open(file_path)
        text = ""
        for page in doc:
            text += page.get_text()
        doc.close()
        
        if not text.strip():
            # If no text found, we'll label it as an image with filename
            text = f"[Image File: {os.path.basename(file_path)} - No text extracted]"
            
        return _chunk_text([{"page": 1, "text": text}], 1)
    except:
        return _chunk_text([{"page": 1, "text": f"[Image: {os.path.basename(file_path)}]" }], 1)


def _chunk_text(pages_text: List[Dict], total_pages: int) -> Tuple[List[Dict], int]:
    if not pages_text:
        return [], total_pages

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len,
        separators=["\n\n", "\n", " ", ""],
    )

    chunks = []
    for page_info in pages_text:
        page_chunks = splitter.create_documents(
            [page_info["text"]],
            metadatas=[{"page": page_info["page"]}],
        )
        for chunk in page_chunks:
            if chunk.page_content.strip():
                chunks.append({
                    "content": chunk.page_content.strip(),
                    "page": page_info["page"],
                })

    return chunks, total_pages
